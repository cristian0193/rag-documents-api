"""Shared pytest fixtures for all tests."""
import uuid
import pytest
from unittest.mock import AsyncMock, MagicMock
from httpx import AsyncClient, ASGITransport
from sqlalchemy.ext.asyncio import create_async_engine, async_sessionmaker

# ── SQLite in-memory database ──────────────────────────────────────────────

TEST_DATABASE_URL = "sqlite+aiosqlite:///:memory:"


@pytest.fixture(scope="function")
async def db_engine():
    """Create an async in-memory SQLite engine with all tables."""
    engine = create_async_engine(TEST_DATABASE_URL, echo=False)
    async with engine.begin() as conn:
        from rag.db.models import Base
        await conn.run_sync(Base.metadata.create_all)
    yield engine
    async with engine.begin() as conn:
        from rag.db.models import Base
        await conn.run_sync(Base.metadata.drop_all)
    await engine.dispose()


@pytest.fixture(scope="function")
async def db_session(db_engine):
    """Provide an async SQLite session for tests."""
    session_factory = async_sessionmaker(db_engine, expire_on_commit=False)
    async with session_factory() as session:
        yield session


# ── Mock ChromaDB client ────────────────────────────────────────────────────

@pytest.fixture
def mock_chroma():
    """Mock ChromaDB client with sensible defaults."""
    chroma = MagicMock()
    chroma.upsert = AsyncMock()
    chroma.query = AsyncMock(return_value={
        "ids": [["doc1_0"]],
        "documents": [["Sample chunk text"]],
        "metadatas": [[{
            "document_id": str(uuid.uuid4()),
            "filename": "test.txt",
            "chunk_index": 0,
        }]],
        "distances": [[0.1]],
    })
    chroma.delete_by_ids = AsyncMock()
    chroma.count = AsyncMock(return_value=1)
    return chroma


# ── Mock Embedder ───────────────────────────────────────────────────────────

@pytest.fixture
def mock_embedder():
    """Mock Embedder that returns the correct number of 384-dim vectors."""
    embedder = MagicMock()

    async def embed_side_effect(texts):
        return [[0.1] * 384 for _ in texts]

    embedder.embed = AsyncMock(side_effect=embed_side_effect)
    embedder.embed_query = AsyncMock(return_value=[0.1] * 384)
    return embedder


# ── Mock OllamaClient ──────────────────────────────────────────────────────

@pytest.fixture
def mock_llm():
    """Mock OllamaClient with canned responses."""
    llm = MagicMock()
    llm.generate = AsyncMock(return_value="This is the generated answer.")
    llm.health_check = AsyncMock(return_value=True)
    llm.build_rag_prompt = MagicMock(return_value="Full RAG prompt here")
    return llm


# ── Sample files ────────────────────────────────────────────────────────────

@pytest.fixture
def sample_docx_bytes() -> bytes:
    """Minimal .docx with one paragraph ('Hello from DOCX')."""
    import io
    import docx as python_docx
    doc = python_docx.Document()
    doc.add_paragraph("Hello from DOCX")
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_xlsx_bytes() -> bytes:
    """Minimal .xlsx with Sheet1 containing 'Hello from XLSX'."""
    import io
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    ws["A1"] = "Hello from XLSX"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_bytes() -> bytes:
    """Minimal .pptx with one slide containing 'Hello from PPTX'."""
    import io
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide_layout = prs.slide_layouts[5]  # blank layout
    slide = prs.slides.add_slide(slide_layout)
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    txBox.text_frame.text = "Hello from PPTX"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_txt_bytes() -> bytes:
    return (
        b"This is a sample text document for testing. "
        b"It has multiple sentences. The content is about testing."
    )


@pytest.fixture
def sample_pdf_bytes() -> bytes:
    """Minimal valid PDF bytes with extractable text ('Hello from PDF')."""
    return b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj

2 0 obj
<< /Type /Pages /Kids [3 0 R] /Count 1 >>
endobj

3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj

4 0 obj
<< /Length 44 >>
stream
BT /F1 12 Tf 100 700 Td (Hello from PDF) Tj ET
endstream
endobj

5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj

xref
0 6
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000360 00000 n

trailer
<< /Size 6 /Root 1 0 R >>
startxref
441
%%EOF"""


# ── FastAPI test client with all overrides ─────────────────────────────────

@pytest.fixture
async def client(db_session, mock_chroma, mock_embedder, mock_llm):
    """AsyncClient with all external services mocked and DB backed by SQLite."""
    from rag.main import app
    from rag.db.postgres import get_db
    from rag.api.deps import get_chroma
    from rag.services.ingestion import IngestionService
    from rag.services.retrieval import RetrievalService

    async def override_get_db():
        yield db_session

    def override_get_chroma():
        return mock_chroma

    app.dependency_overrides[get_db] = override_get_db
    app.dependency_overrides[get_chroma] = override_get_chroma

    # Set app.state manually — lifespan does not run in tests
    app.state.chroma = mock_chroma
    app.state.embedder = mock_embedder
    app.state.llm = mock_llm
    app.state.ingestion_service = IngestionService(
        embedder=mock_embedder,
        chroma=mock_chroma,
    )
    app.state.retrieval_service = RetrievalService(
        embedder=mock_embedder,
        chroma=mock_chroma,
        llm=mock_llm,
    )

    async with AsyncClient(
        transport=ASGITransport(app=app),
        base_url="http://test",
    ) as ac:
        yield ac

    app.dependency_overrides.clear()
    for attr in ("chroma", "embedder", "llm", "ingestion_service", "retrieval_service"):
        if hasattr(app.state, attr):
            delattr(app.state, attr)
