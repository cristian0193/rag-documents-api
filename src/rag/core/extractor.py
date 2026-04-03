"""Text extraction from PDF, TXT, DOCX, XLSX, and PPTX files."""
import io
import re
from pathlib import Path

import docx
import openpyxl
import pypdf
from pptx import Presentation


class UnsupportedFileTypeError(ValueError):
    pass


def _extract_txt(file_bytes: bytes, filename: str) -> str:
    try:
        return file_bytes.decode("utf-8")
    except UnicodeDecodeError:
        return file_bytes.decode("latin-1")


def _extract_pdf(file_bytes: bytes, filename: str) -> str:
    try:
        reader = pypdf.PdfReader(io.BytesIO(file_bytes))
    except pypdf.errors.PdfReadError as e:
        raise ValueError(f"Failed to read PDF '{filename}': {e}") from e
    if len(reader.pages) == 0:
        raise ValueError(f"PDF '{filename}' is empty or unreadable.")
    pages_text = [p.extract_text() for p in reader.pages if p.extract_text()]
    text = "\n".join(pages_text)
    if not text.strip():
        raise ValueError(
            f"Could not extract text from PDF '{filename}'. The file may be image-only or encrypted."
        )
    return text


def _extract_docx(file_bytes: bytes, filename: str) -> str:
    doc = docx.Document(io.BytesIO(file_bytes))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    if not paragraphs:
        raise ValueError(f"Could not extract text from '{filename}'. The document may be empty.")
    return "\n".join(paragraphs)


def _extract_xlsx(file_bytes: bytes, filename: str) -> str:
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            row_text = "\t".join(str(cell) for cell in row if cell is not None)
            if row_text.strip():
                parts.append(row_text)
    if not parts:
        raise ValueError(f"Could not extract text from '{filename}'. The workbook may be empty.")
    return "\n".join(parts)


def _extract_pptx(file_bytes: bytes, filename: str) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    parts = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_texts.append(t)
        if slide_texts:
            parts.append(f"Slide {i}:")
            parts.extend(slide_texts)
    if not parts:
        raise ValueError(
            f"Could not extract text from '{filename}'. The presentation may be empty."
        )
    return "\n".join(parts)


_LEGACY_HINTS = {
    ".doc": "Please convert to .docx before uploading.",
    ".xls": "Please convert to .xlsx before uploading.",
    ".ppt": "Please convert to .pptx before uploading.",
}

_EXTRACTORS = {
    ".txt": _extract_txt,
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".xlsx": _extract_xlsx,
    ".pptx": _extract_pptx,
}


def validate_file_type(filename: str) -> None:
    """
    Raise UnsupportedFileTypeError early (before any DB writes) if the
    file extension is not supported or is a legacy Office format.
    """
    suffix = Path(filename).suffix.lower()
    if suffix in _LEGACY_HINTS:
        raise UnsupportedFileTypeError(
            f"Unsupported file type '{suffix}'. {_LEGACY_HINTS[suffix]}"
        )
    if suffix not in _EXTRACTORS:
        raise UnsupportedFileTypeError(
            f"Unsupported file type '{suffix}'. Supported: .pdf, .txt, .docx, .xlsx, .pptx"
        )


def extract_text(file_bytes: bytes, filename: str) -> str:
    """
    Extract raw text from file bytes.

    Args:
        file_bytes: Raw file content
        filename: Original filename (used to determine file type)

    Returns:
        Extracted text as a string

    Raises:
        UnsupportedFileTypeError: If file type is not supported
        ValueError: If file is empty or unreadable
    """
    suffix = Path(filename).suffix.lower()

    if suffix in _LEGACY_HINTS:
        raise UnsupportedFileTypeError(
            f"Unsupported file type '{suffix}'. {_LEGACY_HINTS[suffix]}"
        )

    extractor = _EXTRACTORS.get(suffix)
    if extractor is None:
        raise UnsupportedFileTypeError(
            f"Unsupported file type '{suffix}'. Supported: .pdf, .txt, .docx, .xlsx, .pptx"
        )

    text = extractor(file_bytes, filename)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()
